#!/usr/bin/env python
# Copyright 2026 Bytedance Ltd. and/or its affiliates
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""Async-HPT 발표 덱 빌더 (30분 / 본편 19장 + Backup).

디자인 언어는 AsyncHPT_Figure.pptx를 계승한다:
  20 x 11.25 in 캔버스 · Noto Sans KR(본문) + Georgia(수식)
  slate 잉크(#0F172A/#334155/#5B677D) · RL=파랑 · SFT τ*=주황 · 정합=teal · 유휴=핑크

사용법: (ppt conda 환경에서)  python make_asynchpt_slides.py
출력:   AsyncHPT_Slides.pptx (스크립트와 같은 디렉터리)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------- palette ----------------
INK    = RGBColor(0x0F, 0x17, 0x2A)
SLATE  = RGBColor(0x33, 0x41, 0x55)
MUTE   = RGBColor(0x5B, 0x67, 0x7D)
LINE   = RGBColor(0xCB, 0xD5, 0xE1)
PANEL  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_BG   = RGBColor(0xDB, 0xEA, 0xFE)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
ORANGE_BG = RGBColor(0xFF, 0xED, 0xD5)
TEAL   = RGBColor(0x0F, 0x76, 0x6E)
TEAL_BG   = RGBColor(0xE6, 0xFF, 0xFA)
RED    = RGBColor(0xB9, 0x1C, 0x1C)
RED_BG    = RGBColor(0xFE, 0xE2, 0xE2)
DARK   = RGBColor(0x1E, 0x29, 0x3B)
GEN    = RGBColor(0xC9, 0xD2, 0xDE)

BODY_F = "Noto Sans KR"
MATH_F = "Georgia"
SCALE = 1.15  # 전역 글자 크기 배율

prs = Presentation()
prs.slide_width = Inches(20)
prs.slide_height = Inches(11.25)
BLANK = prs.slide_layouts[6]
PAGE = [0]

# ---------------- helpers ----------------

def _fonts(run, name):
    run.font.name = name
    rPr = run.font._rPr
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def R(t, **kw):
    d = {'t': t}
    d.update(kw)
    return d


def M(t, **kw):
    kw.setdefault('f', MATH_F)
    kw.setdefault('i', True)
    return R(t, **kw)


def _fill_runs(p, runs_, size, color, bold, font):
    for r in runs_:
        run = p.add_run()
        run.text = r['t']
        f = run.font
        _fonts(run, r.get('f', font))
        f.size = Pt(round(r.get('s', size) * SCALE, 1))
        f.bold = r.get('b', bold)
        f.italic = r.get('i', False)
        f.color.rgb = r.get('c', color)
        if r.get('sub'):
            f._rPr.set('baseline', '-25000')
        if r.get('sup'):
            f._rPr.set('baseline', '30000')


def tb(sl, x, y, w, h, lines, size=17, color=SLATE, bold=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.18,
       font=BODY_F, space_after=5):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = leading
        p.space_after = Pt(space_after)
        if isinstance(ln, str):
            ln = [R(ln)]
        _fill_runs(p, ln, size, color, bold, font)
    return box


def panel(sl, x, y, w, h, fill=PANEL, line=LINE, lw=1.2, radius=0.055, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def dash_border(sp, val='dash'):
    ln = sp.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': val}))


def chip(sl, x, y, w, h, text=None, fill=BLUE, tcolor=WHITE, size=13, bold=True,
         line=None, radius=0.28, font=BODY_F, runs_=None):
    sp = panel(sl, x, y, w, h, fill=fill, line=line, radius=radius)
    if text is not None or runs_ is not None:
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _fill_runs(p, runs_ if runs_ is not None else [R(text)], size, tcolor, bold, font)
    return sp


def seg(sl, x1, y1, x2, y2, color=SLATE, w=1.5, dash=None, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    c.shadow.inherit = False
    return c


def add_slide(kicker=None, title=None, footer=True):
    sl = prs.slides.add_slide(BLANK)
    PAGE[0] += 1
    if kicker:
        tb(sl, 0.9, 0.52, 18.2, 0.4, [[R(kicker, b=True, c=TEAL, s=15)]])
    if title:
        tb(sl, 0.9, 0.9, 18.2, 0.75, [[R(title, b=True, c=INK, s=30)]])
        seg(sl, 0.9, 1.78, 19.1, 1.78, color=LINE, w=1.2)
    if footer:
        tb(sl, 0.9, 10.74, 12, 0.35, [[R('Async-HPT — 의미 보존형 비동기 Hybrid Post-Training', c=MUTE, s=11.5)]])
        tb(sl, 17.9, 10.74, 1.2, 0.35, [[R(str(PAGE[0]), c=MUTE, s=11.5)]], align=PP_ALIGN.RIGHT)
    return sl


def legend(sl, x, y, items, size=12.5):
    cx = x
    for label, fill, border in items:
        panel(sl, cx, y, 0.28, 0.28, fill=fill, line=border, radius=0.3)
        t = tb(sl, cx + 0.38, y - 0.035, 1.9, 0.35, [[R(label, c=SLATE, s=size, b=True)]])
        cx += 0.38 + 0.16 * len(label) + 0.75


def band(sl, x, y, w, h, lines, fill=TEAL_BG, border=TEAL, tcolor=None, size=16.5, bold=False, align=PP_ALIGN.LEFT):
    panel(sl, x, y, w, h, fill=fill, line=border, lw=1.4, radius=0.10)
    tb(sl, x + 0.45, y, w - 0.9, h, lines, size=size, color=tcolor or SLATE, bold=bold,
       align=align, anchor=MSO_ANCHOR.MIDDLE)


# ================= S1 · 타이틀 =================
sl = add_slide(footer=False)
tb(sl, 0.9, 3.0, 18.2, 0.5, [[R('논문 발표', b=True, c=TEAL, s=17)]], align=PP_ALIGN.CENTER)
tb(sl, 0.9, 3.6, 18.2, 1.3, [[R('Async-HPT', b=True, c=INK, s=64)]], align=PP_ALIGN.CENTER)
tb(sl, 0.9, 5.05, 18.2, 0.7, [[R('비동기 experience stream 위의 의미 보존형 Hybrid Post-Training', c=SLATE, s=26)]],
   align=PP_ALIGN.CENTER)
seg(sl, 7.5, 6.15, 12.5, 6.15, color=LINE, w=1.5)
tb(sl, 0.9, 6.45, 18.2, 0.6,
   [[R('두 학습 신호가 하나의 비동기 stream을 공유하면서도, 각자의 estimator 의미를 온전히 보존한다',
       c=MUTE, s=18, i=True)]], align=PP_ALIGN.CENTER)
tb(sl, 0.9, 8.9, 18.2, 0.9,
   [[R('발표자 · 서강대학교', c=SLATE, s=17, b=True)],
    [R('2026. 7.', c=MUTE, s=14)]], align=PP_ALIGN.CENTER)

# ================= S2 · 문제의식 ① =================
sl = add_slide('문제의식 ①', 'Agent 시대의 post-training — 두 부담이 함께 온다')
panel(sl, 0.9, 2.15, 8.85, 5.75)
tb(sl, 1.35, 2.55, 8.0, 0.6, [[R('시간의 부담', b=True, c=BLUE, s=21)]])
tb(sl, 1.35, 3.3, 8.0, 0.8, [[R('“생성이 학습을 기다리게 한다”', b=True, c=INK, s=23)]])
turns = [('관찰', GEN), ('추론·행동', DARK), ('관찰', GEN), ('추론·행동', DARK), ('…', PANEL), ('종료', DARK)]
xx = 1.35
for t, f in turns:
    ww = 0.55 + 0.155 * len(t)
    chip(sl, xx, 4.45, ww, 0.62, t, fill=f, tcolor=(SLATE if f in (GEN, PANEL) else WHITE), size=12,
         line=LINE if f is PANEL else None, radius=0.22)
    xx += ww + 0.12
tb(sl, 1.35, 5.55, 7.95, 2.2,
   [[R('computer-use trajectory는 수십 turn의 관찰·행동 상호작용 —', c=SLATE, s=16)],
    [R('generation 시간이 학습 시간을 압도하고, 길이 편차가 크다.', c=SLATE, s=16)],
    [R('가장 긴 rollout 하나가 전체를 기다리게 한다.', c=SLATE, s=16)]])
panel(sl, 10.25, 2.15, 8.85, 5.75)
tb(sl, 10.7, 2.55, 8.0, 0.6, [[R('신호의 부담', b=True, c=ORANGE, s=21)]])
tb(sl, 10.7, 3.3, 8.2, 0.8, [[R('“성공이 없으면 gradient도 없다”', b=True, c=INK, s=23)]])
for k in range(6):
    chip(sl, 10.7 + k * 1.05, 4.45, 0.9, 0.62, '✗', fill=RED_BG, tcolor=RED, size=15, line=RED, radius=0.22)
tb(sl, 10.7, 5.55, 8.1, 2.2,
   [[R('검증 가능한 보상(RLVR)은 명확하지만, 어려운 과제에서는', c=SLATE, s=16)],
    [R('성공 자체가 드물다.', c=SLATE, s=16)],
    [R('실패만 있는 rollout group은 아무것도 가르치지 못한다.', c=SLATE, s=16)]])
band(sl, 0.9, 8.45, 18.2, 1.35,
     [[R('같은 post-training 시스템이 두 부담을 동시에 진다 — ', b=True, c=TEAL, s=20),
       R('하나만 푸는 해법으로는 부족하다.', c=SLATE, s=20)]])

# ================= S3 · 문제의식 ② =================
sl = add_slide('문제의식 ②', '각자의 해법은 성숙했다 — 함께 쓸 수 있는가?')
panel(sl, 0.9, 2.15, 8.0, 2.5, fill=BLUE_BG, line=BLUE, lw=1.6)
tb(sl, 1.35, 2.45, 7.2, 0.55, [[R('효율의 해법 — 비동기 RL', b=True, c=BLUE, s=20)]])
tb(sl, 1.35, 3.2, 7.2, 1.2,
   [[R('생성과 학습을 stream으로 겹쳐 유휴를 없앤다', c=INK, s=17, b=True)],
    [R('StreamRL · AsyncFlow …', c=MUTE, s=14)]])
panel(sl, 11.1, 2.15, 8.0, 2.5, fill=ORANGE_BG, line=ORANGE, lw=1.6)
tb(sl, 11.55, 2.45, 7.2, 0.55, [[R('신호의 해법 — Hybrid Post-Training', b=True, c=ORANGE, s=20)]])
tb(sl, 11.55, 3.2, 7.2, 1.2,
   [[R('실패 영역에 검증된 trajectory의 supervised 신호를 주입한다', c=INK, s=17, b=True)],
    [R('LUFFY · UPT/HPT …', c=MUTE, s=14)]])
seg(sl, 4.9, 4.65, 8.6, 5.75, color=BLUE, w=2.2, arrow=True)
seg(sl, 15.1, 4.65, 11.4, 5.75, color=ORANGE, w=2.2, arrow=True)
chip(sl, 8.15, 5.85, 3.7, 1.15, '하나의 learner', fill=INK, size=19, radius=0.14)
band(sl, 0.9, 7.55, 18.2, 1.4,
     [[R('이 발표의 질문 — ', b=True, c=TEAL, s=21),
       R('두 해법을 하나의 learner에 결합할 때, 각자가 서 있는 전제는 그대로 성립하는가?', c=INK, s=21)]])
tb(sl, 0.9, 9.3, 18.2, 0.6,
   [[R('답부터: 그대로 포개면 성립하지 않는다 — 그래서 ', c=SLATE, s=17),
     R('의미 보존형 결합', c=TEAL, s=17, b=True), R('이 필요하다.', c=SLATE, s=17)]],
   align=PP_ALIGN.CENTER)

# ================= S4 · 배경 ① 실행의 축 (타임라인 그림) =================
sl = add_slide('배경 ①  ·  실행의 축', '동기식 RL은 긴 trajectory 앞에서 멈춘다')
legend(sl, 14.4, 2.05, [('생성', GEN, None), ('학습', DARK, None), ('유휴', RED_BG, RED)])

# sync row
tb(sl, 0.9, 3.15, 1.5, 0.5, [[R('sync', b=True, c=INK, s=19)]])
x = 2.6; y0 = 2.9; hh = 0.8
sync_blocks = [(2.2, GEN, '생성'), (0.6, RED_BG, '유휴'), (1.5, DARK, '학습'), (0.6, RED_BG, '유휴'),
               (3.4, GEN, '생성 · long-tail'), (0.9, RED_BG, '유휴'), (1.5, DARK, '학습'),
               (0.6, RED_BG, '유휴'), (2.2, GEN, '생성')]
for w, f, lab in sync_blocks:
    tc = WHITE if f == DARK else (RED if f == RED_BG else SLATE)
    bd = RED if f == RED_BG else None
    chip(sl, x, y0, w, hh, lab, fill=f, tcolor=tc, size=13, bold=True, line=bd, radius=0.14)
    x += w + 0.06
sync_end = x - 0.06
tb(sl, 2.6, 3.95, 15.5, 0.4,
   [[R('생성과 학습이 ', c=MUTE, s=14), R('교대', c=INK, s=14, b=True),
     R(' → 사이마다 ', c=MUTE, s=14), R('유휴(bubble)', c=RED, s=14, b=True),
     R(' · 소수의 long-tail rollout이 전체 generation을 지연', c=MUTE, s=14)]])

# async rows
tb(sl, 0.9, 6.0, 1.6, 0.5, [[R('async', b=True, c=INK, s=19)]])
tb(sl, 2.6, 5.28, 4.0, 0.35, [[R('생성 (streaming)', c=MUTE, s=12.5, b=True)]])
x = 2.6
for i in range(12):
    chip(sl, x, 5.62, 0.62, 0.52, fill=GEN, radius=0.16)
    x += 0.70
gen_end = x - 0.08
x = 3.15
for i in range(12):
    chip(sl, x, 6.28, 0.58, 0.52, fill=DARK, radius=0.16)
    x += 0.66
tr_end = x - 0.08
async_end = max(gen_end, tr_end)
tb(sl, 2.6, 6.92, 4.0, 0.35, [[R('학습 (streaming)', c=MUTE, s=12.5, b=True)]])
tb(sl, 2.6, 7.38, 15.5, 0.4,
   [[R('생성과 학습이 ', c=MUTE, s=14), R('겹침', c=INK, s=14, b=True),
     R(' — 유휴 제거 · long-tail은 attempt 독립 실행과 queue가 흡수', c=MUTE, s=14)]])

# time axis + saved-time bracket
seg(sl, 2.6, 8.25, 18.3, 8.25, color=SLATE, w=1.6, arrow=True)
tb(sl, 18.35, 8.06, 0.75, 0.4, [[R('시간', c=SLATE, s=13, b=True)]])
seg(sl, async_end, 8.75, sync_end, 8.75, color=TEAL, w=2.2)
seg(sl, async_end, 8.6, async_end, 8.9, color=TEAL, w=2.2)
seg(sl, sync_end, 8.6, sync_end, 8.9, color=TEAL, w=2.2)
tb(sl, async_end, 8.98, sync_end - async_end, 0.4, [[R('단축된 벽시계 시간', c=TEAL, s=14, b=True)]],
   align=PP_ALIGN.CENTER)

band(sl, 0.9, 9.6, 18.2, 0.95,
     [[R('비동기 RL 계보 — StreamRL · AsyncFlow. ', c=SLATE, b=True),
       R('generation과 training의 의존성을 stream으로 끊는다.  대가: 학습 sample이 과거 policy의 산물이 된다 (→ 사전지식 ①).', c=SLATE)]],
     fill=PANEL, border=LINE)

# ================= S3 · 배경 ② 신호의 축 =================
sl = add_slide('배경 ②  ·  신호의 축', '어려운 prompt에서는 online rollout만으로 gradient가 없다')
# left: sparse reward
panel(sl, 0.9, 2.1, 8.85, 6.35)
tb(sl, 1.35, 2.4, 8.0, 0.5, [[R('어려운 prompt의 rollout group', b=True, c=INK, s=19)]])
chip(sl, 1.35, 3.35, 1.1, 0.85, runs_=[M('x', s=17), M('i', s=12, sub=True)], fill=BLUE)
tb(sl, 1.35, 4.3, 1.2, 0.4, [[R('× n', c=MUTE, s=13, b=True)]], align=PP_ALIGN.CENTER)
for k in range(8):
    yy = 2.95 + k * 0.54
    chip(sl, 3.1, yy, 3.3, 0.42, fill=WHITE, line=LINE, radius=0.2)
    tb(sl, 6.6, yy - 0.02, 1.7, 0.45, [[R('✗  r = 0', c=RED, s=14, b=True)]])
seg(sl, 2.55, 3.75, 3.0, 3.75, color=MUTE, w=1.4, arrow=True)
tb(sl, 1.35, 7.55, 7.9, 0.75,
   [[R('전 rollout 실패 → group 안에서 차이가 없음 → ', c=SLATE, s=15.5),
     R('advantage = 0, gradient 없음', c=RED, s=15.5, b=True)]])
# right: remedy
panel(sl, 10.25, 2.1, 8.85, 6.35)
tb(sl, 10.7, 2.4, 8.0, 0.5, [[R('해법 — 검증된 trajectory의 supervised 신호', b=True, c=INK, s=19)]])
chip(sl, 10.7, 3.35, 1.3, 0.85, runs_=[M('τ', s=17), M('★', s=12, sup=True)], fill=ORANGE)
tb(sl, 12.3, 3.42, 6.6, 0.8,
   [[R('같은 prompt에 대응하는 verified trajectory —', c=SLATE, s=15.5)],
    [R('실패 영역에는 supervised 신호를, 그 외에는 online RL 신호를', c=SLATE, s=15.5)]])
tb(sl, 10.7, 4.9, 8.0, 0.4, [[R('계보', b=True, c=TEAL, s=15)]])
tb(sl, 10.7, 5.35, 8.1, 2.6,
   [[R('LUFFY', c=INK, s=15.5, b=True), R(' — off-policy guidance: 실패 영역에 외부 trajectory 주입', c=SLATE, s=15.5)],
    [R('UPT / HPT', c=INK, s=15.5, b=True),
     R(' — SFT와 RL을 하나의 estimator로 통일, 성능 피드백으로 전환 (→ 사전지식 ②)', c=SLATE, s=15.5)]])
band(sl, 0.9, 8.8, 18.2, 0.95,
     [[R('RLVR — 규칙 기반 verifier, 보상 r ∈ [0, 1]. ', c=SLATE, b=True),
       R('명시적이고 재현 가능하지만, 난이도가 높은 영역에서 보상이 희박해진다.', c=SLATE)]],
     fill=PANEL, border=LINE)

# ================= S4 · 충돌 =================
sl = add_slide('문제', '두 축은 learner의 경계에서 충돌한다')
panel(sl, 0.9, 2.1, 8.85, 4.2, fill=BLUE_BG, line=BLUE, lw=1.6)
tb(sl, 1.35, 2.45, 8.0, 0.5, [[R('비동기 축의 전제', b=True, c=BLUE, s=19)]])
tb(sl, 1.35, 3.15, 8.0, 1.0, [[R('“모든 sample은 rollout policy의 산물이다”', b=True, c=INK, s=23)]])
tb(sl, 1.35, 4.35, 8.0, 1.6,
   [[R('생성/학습 시점의 policy 차이를 보정하고(correction), 지나치게 오래된 sample을', c=SLATE, s=15.5)],
    [R('거르는(staleness) 모든 장치가 이 전제 위에 서 있다.', c=SLATE, s=15.5)]])
panel(sl, 10.25, 2.1, 8.85, 4.2, fill=ORANGE_BG, line=ORANGE, lw=1.6)
tb(sl, 10.7, 2.45, 8.0, 0.5, [[R('hybrid 축의 주입물', b=True, c=ORANGE, s=19)]])
tb(sl, 10.7, 3.15, 8.2, 1.0, [[R('“verified trajectory τ★는 rollout이 아니라 고정 target이다”', b=True, c=INK, s=22)]])
tb(sl, 10.7, 4.35, 8.2, 1.8,
   [[R('생성 시점의 policy도, rollout log-probability', c=SLATE, s=15.5),
     R('(생성 당시 “얼마나 확신했나”의 기록)', c=MUTE, s=13.5), R('도 없다.', c=SLATE, s=15.5)],
    [R('비동기 기제가 요구하는 provenance(출처 증명) 자체가 존재하지 않는다.', c=SLATE, s=15.5)]])
seg(sl, 9.75, 4.2, 10.25, 4.2, color=RED, w=2.5)
band(sl, 0.9, 6.75, 18.2, 1.7,
     [[R('그대로 포개면 — supervised 신호가 자신을 위해 만들어지지 않은 보정과 필터에 노출된다.', c=RED, b=True, s=19)],
      [R('예: placeholder log-probability가 policy mismatch의 “증거”로 읽혀, τ★가 mask · reweight · 폐기될 수 있다.', c=SLATE, s=15.5)]],
     fill=RED_BG, border=RED)
tb(sl, 0.9, 8.85, 18.2, 0.6,
   [[R('효율을 위한 asynchrony와 신호를 위한 hybrid estimator는 경계에서 충돌한다.', b=True, c=INK, s=19)]],
   align=PP_ALIGN.CENTER)

# ================= S5 · 이 논문 =================
sl = add_slide('이 논문', '경계를 설계의 대상으로 — Async-HPT')
band(sl, 0.9, 2.05, 18.2, 1.55,
     [[R('transport와 semantics를 분리한다 — ', b=True, c=TEAL, s=20),
       R('두 신호가 같은 비동기 stream을 공유하면서도, 각자의 estimator 의미를 온전히 보존하는 learner contract.',
         c=SLATE, s=20)]])
cards = [
    ('① 원리', '경험을 운반하는 통로(queue)와 의미가 규정되는 지점(learner boundary)을 분리 — 의미는 운반 중에 해석되지 않는다.'),
    ('② 방법', '네 개의 질문에 대한 설계 — 성능 피드백 라우팅 · 통합 objective(advantage/reference 재사용) · provenance 기반 validity · 실행 단위 분리.'),
    ('③ 실증', 'verl 기반 비동기 스택 위 구현 — coupled anchor 대비 비교와 2×2 ablation으로 각 설계 축의 기여를 분리.'),
]
for i, (t, d) in enumerate(cards):
    x = 0.9 + i * 6.2
    panel(sl, x, 4.0, 5.8, 2.9)
    tb(sl, x + 0.4, 4.3, 5.0, 0.5, [[R(t, b=True, c=INK, s=19)]])
    tb(sl, x + 0.4, 4.95, 5.0, 1.8, [[R(d, c=SLATE, s=15)]])
roadmap = ['배경', '원리', 'Q1 라우팅', 'Q2 objective', 'Q3 validity', 'Q4 실행', '실험', '결론']
x = 0.9
tb(sl, 0.9, 7.35, 4.0, 0.4, [[R('로드맵', b=True, c=TEAL, s=15)]])
for i, t in enumerate(roadmap):
    w = 2.05
    chip(sl, x, 7.85, w, 0.62, t, fill=WHITE if i not in (2, 3, 4, 5) else TEAL_BG,
         line=LINE if i not in (2, 3, 4, 5) else TEAL, tcolor=SLATE if i not in (2, 3, 4, 5) else TEAL,
         size=13.5, radius=0.3)
    if i < len(roadmap) - 1:
        seg(sl, x + w + 0.04, 8.16, x + w + 0.24, 8.16, color=MUTE, w=1.3, arrow=True)
    x += w + 0.28

# ================= S6 · 사전지식 ① validity 기제 =================
sl = add_slide('사전지식 ①', '비동기 RL의 validity 기제 — correction과 staleness')
for i, (t, d1, d2, d3) in enumerate([
    ('Off-policy correction',
     [R('sample은 생성 시점 policy ', c=SLATE), M('π', c=SLATE), M('g', c=SLATE, sub=True),
      R('의 산물, 학습은 그 뒤의 ', c=SLATE), M('π', c=SLATE), M('θ', c=SLATE, sub=True), R('에서.', c=SLATE)],
     [R('ratio ', c=SLATE), M('π', c=SLATE), M('θ', c=SLATE, sub=True), M('/π', c=SLATE),
      M('g', c=SLATE, sub=True),
      R('로 mismatch를 보정하고, 지나친 어긋남은 절단·기각한다 (truncated IS 등).', c=SLATE)],
     [R('직관 — ', c=TEAL, s=14, b=True),
      R('생성 당시보다 확률이 많이 달라진 token일수록, 그 경험을 그만큼 덜 믿는다.', c=MUTE, s=14)]),
    ('Staleness filtering',
     [R('sample마다 생성 version ', c=SLATE), M('g', c=SLATE), R('가 기록된다.', c=SLATE)],
     [R('학습 version과의 격차가 한도를 넘으면 — 너무 오래된 sample — 학습에서 제외한다.', c=SLATE)],
     [R('직관 — ', c=TEAL, s=14, b=True),
      R('너무 옛날 policy의 경험은 지금의 나와 무관한 이야기가 되기 때문이다.', c=MUTE, s=14)]),
]):
    x = 0.9 + i * 9.35
    panel(sl, x, 2.1, 8.85, 5.05)
    tb(sl, x + 0.45, 2.45, 8.0, 0.5, [[R(t, b=True, c=INK, s=20)]])
    # mini timeline
    yy = 3.35
    for k in range(5):
        f = BLUE if k == 0 else (DARK if k == 4 else GEN)
        lab = [M('π', c=WHITE, s=13), M('g', c=WHITE, s=10, sub=True)] if k == 0 else \
              ([M('π', c=WHITE, s=13), M('θ', c=WHITE, s=10, sub=True)] if k == 4 else None)
        chip(sl, x + 0.45 + k * 1.65, yy, 1.35, 0.6, runs_=lab, fill=f, radius=0.2)
    seg(sl, x + 0.45, yy + 0.95, x + 0.45 + 4 * 1.65 + 1.35, yy + 0.95, color=MUTE, w=1.3, arrow=True)
    tb(sl, x + 0.45, yy + 1.05, 7.9, 0.35, [[R('policy version →', c=MUTE, s=12)]])
    tb(sl, x + 0.45, 5.05, 7.95, 2.0, [d1, d2, d3], size=15.5)
band(sl, 0.9, 7.55, 18.2, 1.5,
     [[R('두 기제 모두 provenance에 의존한다 — ', b=True, c=TEAL, s=19),
       R('이 sample이 어느 policy에서, 언제 생성되었는가.', c=SLATE, s=19)],
      [R('rollout log-probability와 생성 version이 그 증거다.  (이 문장은 Q3에서 회수된다)', c=MUTE, s=15)]])

# ================= S7 · 사전지식 ② unified 관점 =================
sl = add_slide('사전지식 ②', 'Unified 관점 — SFT와 RL은 같은 estimator의 다른 슬롯 선택')
tb(sl, 0.9, 2.0, 18.2, 0.5,
   [[R('UPT: ', b=True, c=INK, s=15.5),
     R('SFT와 RL은 별개의 목적이 아니라, 하나의 policy-gradient estimator에서 네 슬롯의 선택이 다를 뿐이다.  ', c=SLATE, s=15.5),
     R('HPT: ', b=True, c=INK, s=15.5),
     R('rollout 성능 피드백으로 prompt마다 슬롯을 전환한다.', c=SLATE, s=15.5)]])
band(sl, 2.4, 2.72, 15.36, 0.78,
     [[R('읽는 눈 — ', b=True, c=TEAL, s=14.5),
       R('모든 policy-gradient 학습 신호는  ', c=SLATE, s=14.5),
       R('advantage', b=True, c=INK, s=14.5), R('(얼마나 좋았나)', c=MUTE, s=13),
       R(' × ', c=SLATE, s=14.5),
       R('ratio', b=True, c=INK, s=14.5), R('(그 행동의 확률을, 어느 기준 — 분모 — 대비 키울까)', c=MUTE, s=13),
       R('의 곱이다.', c=SLATE, s=14.5)]])
# 4-slot table
tx, ty = 2.4, 3.75
col_w = [4.0, 5.6, 5.6]
row_h = 0.95
headers = ['슬롯', 'RL', 'SFT']
rows = [
    ('data source', [M('τ'), R(' ~ rollout  '), M('π'), M('g', sub=True)], [R('verified trajectory '), M('τ'), M('★', sup=True)]),
    ('reference (분모)', [R('rollout policy  '), M('π'), M('g', sub=True)], [R('current policy  '), M('π'), M('θ', sub=True)]),
    ('advantage', [R('group-relative  '), M('Â')], [R('상수 (positive)')]),
    ('stabilization', [R('clip · IS 보정')], [R('없음 ('), M('ρ'), R(' ≡ 1)')]),
]
hx = tx
for j, h in enumerate(headers):
    fill = INK if j == 0 else (BLUE if j == 1 else ORANGE)
    chip(sl, hx, ty, col_w[j], 0.75, h, fill=fill, size=17, radius=0.12)
    hx += col_w[j] + 0.08
for i, (name, rl, sft) in enumerate(rows):
    yy = ty + 0.83 + i * (row_h + 0.08)
    hx = tx
    chip(sl, hx, yy, col_w[0], row_h, name, fill=PANEL, tcolor=INK, size=15.5, line=LINE, radius=0.12)
    hx += col_w[0] + 0.08
    sp = panel(sl, hx, yy, col_w[1], row_h, fill=BLUE_BG, line=None, radius=0.12)
    tb(sl, hx + 0.3, yy, col_w[1] - 0.6, row_h, [[R(r['t'], **{k: v for k, v in r.items() if k != 't'}) for r in rl]],
       size=15.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    hx += col_w[1] + 0.08
    panel(sl, hx, yy, col_w[2], row_h, fill=ORANGE_BG, line=None, radius=0.12)
    tb(sl, hx + 0.3, yy, col_w[2] - 0.6, row_h, [[R(r['t'], **{k: v for k, v in r.items() if k != 't'}) for r in sft]],
       size=15.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
tb(sl, 0.9, 8.72, 18.2, 0.45,
   [[R('rollout이 성공하는 영역', c=BLUE, s=15, b=True), R(' → RL 슬롯   ·   ', c=SLATE, s=15),
     R('실패하는 영역', c=ORANGE, s=15, b=True),
     R(' → SFT 슬롯   —   두 신호가 하나의 gradient 문법을 공유한다.', c=SLATE, s=15)]],
   align=PP_ALIGN.CENTER)
band(sl, 0.9, 9.3, 18.2, 1.05,
     [[R('남는 질문 — ', b=True, c=TEAL, s=17.5),
       R('이 표를 비동기 experience stream 위에서 성립시키려면?  이것이 방법이 답할 질문이다.', c=SLATE, s=17.5)]])

# ================= S8 · 원리 (아키텍처 그림) =================
sl = add_slide('방법 · 원리', 'Transport와 Semantics의 분리')
legend(sl, 12.6, 2.0, [('RL 분기', BLUE, None), ('SFT τ★ 분기', ORANGE, None), ('정합 지점', TEAL, None)])
# param sync dashed
seg(sl, 14.9, 2.62, 14.9, 2.42, color=MUTE, w=1.4)
seg(sl, 14.9, 2.42, 4.5, 2.42, color=MUTE, w=1.4, dash='dash')
seg(sl, 4.5, 2.42, 4.5, 2.62, color=MUTE, w=1.4, arrow=True)
tb(sl, 6.0, 2.06, 8.0, 0.35, [[R('parameter sync — 버전 격차가 staleness로', c=MUTE, s=12.5)]], align=PP_ALIGN.CENTER)
# panels
panel(sl, 0.9, 2.62, 7.2, 5.5)
panel(sl, 8.45, 2.62, 1.95, 5.5)
panel(sl, 10.75, 2.62, 8.35, 5.5)
tb(sl, 1.25, 2.85, 5.5, 0.5, [[R('①  생성', b=True, c=INK, s=19), R('   Rollouter ↔ Environment', c=MUTE, s=13)]])
tb(sl, 8.7, 2.85, 1.6, 0.9, [[R('②  Queue', b=True, c=INK, s=17)]])
tb(sl, 11.1, 2.85, 7.6, 0.5, [[R('③  통합 Learner', b=True, c=INK, s=19), R('   unified estimator', c=MUTE, s=13)]])
# 생성 내부
chip(sl, 1.25, 4.7, 1.05, 0.85, runs_=[M('x', s=16), M('i', s=11, sub=True)], fill=BLUE)
tb(sl, 1.25, 4.25, 1.1, 0.4, [[R('× n', c=MUTE, s=12.5, b=True)]], align=PP_ALIGN.CENTER)
tb(sl, 2.62, 3.7, 2.4, 0.35, [[R('attempt 독립 실행', c=MUTE, s=12)]])
alens = [1.15, 1.35, 1.9, 1.0]
for k, alw in enumerate(alens):
    f = RED_BG if k == 2 else WHITE
    bd = RED if k == 2 else LINE
    chip(sl, 2.62, 4.1 + k * 0.5, alw, 0.36, fill=f, line=bd, radius=0.22)
tb(sl, 2.62, 6.15, 2.3, 0.7, [[R('길이 제각각 —', c=MUTE, s=11.5)], [R('long-tail', c=RED, s=11.5, b=True)]])
seg(sl, 2.38, 5.1, 2.58, 5.1, color=MUTE, w=1.3, arrow=True)
chip(sl, 4.5, 4.35, 1.95, 1.35, fill=TEAL, radius=0.12,
     runs_=[R('Accumulator', c=WHITE, s=12.5, b=True)])
tb(sl, 4.5, 5.78, 1.95, 0.35, [[R('그룹 복원', c=TEAL, s=12, b=True)]], align=PP_ALIGN.CENTER)
seg(sl, 4.28, 5.02, 4.45, 5.02, color=MUTE, w=1.3, arrow=True)
sp = chip(sl, 6.62, 4.35, 1.35, 1.35, fill=WHITE, line=SLATE, radius=0.12,
          runs_=[R('Routing\nGate', c=INK, s=12.5, b=True)])
tb(sl, 6.62, 5.78, 1.35, 0.4, [[M('P', c=SLATE, s=13), M('i', c=SLATE, s=10, sub=True), R(' vs ', c=SLATE, s=12.5), M('γ', c=SLATE, s=13)]],
   align=PP_ALIGN.CENTER)
seg(sl, 6.47, 5.02, 6.6, 5.02, color=MUTE, w=1.3, arrow=True)
tb(sl, 1.25, 7.35, 6.6, 0.5,
   [[M('P', c=SLATE, s=13.5), M('i', c=SLATE, s=10, sub=True), R(' > ', c=SLATE, s=13.5), M('γ', c=SLATE, s=13.5),
     R(' → RL 그룹    ·    ', c=SLATE, s=13.5),
     M('P', c=SLATE, s=13.5), M('i', c=SLATE, s=10, sub=True), R(' ≤ ', c=SLATE, s=13.5), M('γ', c=SLATE, s=13.5),
     R(' → SFT ', c=SLATE, s=13.5), M('τ', c=SLATE, s=13.5), M('★', c=SLATE, s=10, sup=True),
     R('  (없으면 RL)', c=MUTE, s=13.5)]], align=PP_ALIGN.CENTER)
seg(sl, 8.12, 5.02, 8.42, 5.02, color=SLATE, w=1.8, arrow=True)
# queue 내부
qcols = [BLUE, ORANGE, BLUE, BLUE, ORANGE, BLUE, BLUE]
for k, f in enumerate(qcols):
    chip(sl, 8.75, 3.85 + k * 0.52, 1.35, 0.38, fill=f, radius=0.22)
tb(sl, 8.45, 7.5, 1.95, 0.6, [[R('도착 순서 그대로', c=MUTE, s=10.5)], [R('해석 없이 운반', c=MUTE, s=10.5)]],
   align=PP_ALIGN.CENTER, space_after=1)
seg(sl, 10.42, 5.02, 10.72, 5.02, color=SLATE, w=1.8, arrow=True)
# learner 내부
tb(sl, 11.1, 3.5, 1.9, 0.35, [[R('mixed batch', c=SLATE, s=12.5, b=True)]])
bcols = [BLUE, BLUE, BLUE, ORANGE, BLUE, BLUE]
for k, f in enumerate(bcols):
    chip(sl, 11.1, 3.9 + k * 0.52, 1.5, 0.38, fill=f, radius=0.22)
tb(sl, 10.9, 7.1, 2.0, 0.6, [[R('RL n행 + SFT 1행', c=MUTE, s=11.5)]], align=PP_ALIGN.CENTER)
seg(sl, 12.75, 5.35, 13.05, 5.35, color=MUTE, w=1.4, arrow=True)
panel(sl, 13.1, 3.4, 3.6, 4.1, fill=WHITE, line=TEAL, lw=1.8, radius=0.08)
tb(sl, 13.4, 3.62, 3.05, 0.45, [[R('하나의 estimator', b=True, c=TEAL, s=16)]])
seg(sl, 13.4, 4.15, 16.4, 4.15, color=LINE, w=1.0)
tb(sl, 13.4, 4.3, 3.1, 3.1,
   [[R('① advantage', c=INK, s=13.5, b=True)],
    [R('group ↔ β pseudo-reward', c=SLATE, s=12.5)],
    [R('② reference', c=INK, s=13.5, b=True)],
    [R('rollout logπ ↔ self-detach', c=SLATE, s=12.5)],
    [R('③ validity', c=INK, s=13.5, b=True)],
    [R('correction·staleness는 RL 행에만', c=SLATE, s=12.5)]], space_after=3)
seg(sl, 16.75, 5.35, 17.0, 5.35, color=MUTE, w=1.4, arrow=True)
chip(sl, 17.05, 4.75, 1.75, 1.2, fill=BLUE, radius=0.14,
     runs_=[M('π', c=WHITE, s=19), R(' update', c=WHITE, s=15, b=True)])
# 의미 상태 스트립
labels3 = [('route metadata 부착', 4.5), ('해석 없이 운반', 9.42), ('learner boundary에서 의미 실행', 14.9)]
tb(sl, 0.9, 8.35, 3.2, 0.4, [[R('branch 의미의 상태', b=True, c=TEAL, s=14)]])
for t, cx in labels3:
    tb(sl, cx - 2.4, 8.35, 4.8, 0.4, [[R(t, c=SLATE, s=13.5, b=True)]], align=PP_ALIGN.CENTER)
# Q chips
qm = [('Q1 라우팅 — Gate', 0.9), ('Q2 objective — Estimator', 5.55), ('Q3 validity — Correction · Staleness', 10.2), ('Q4 실행 — Attempt · Accumulator', 14.85)]
for t, x in qm:
    chip(sl, x, 9.05, 4.25, 0.65, t, fill=TEAL_BG, line=TEAL, tcolor=TEAL, size=13.5, radius=0.3)

# ================= S9 · Q1 라우팅 =================
sl = add_slide('방법 · Q1', '어떤 prompt에 어떤 신호를? — 성능 피드백 라우팅')
band(sl, 0.9, 2.0, 18.2, 1.0,
     [[M('P', s=19), M('i', s=13, sub=True), M(' = (1/n) Σ', s=19), M('j', s=13, sub=True),
       M(' 1[ s', s=19), M('i,j', s=13, sub=True), M(' > δ ]', s=19),
       R('    —  n번의 online rollout 성공률 = 이 prompt에 대한 policy의 성능', c=SLATE, s=16.5),
       R('    ( 1[·]: 성공이면 1, 아니면 0 )', c=MUTE, s=13)]],
     fill=PANEL, border=LINE)
# diagram
chip(sl, 1.3, 4.6, 1.1, 0.9, runs_=[M('x', s=16), M('i', s=11, sub=True)], fill=BLUE)
seg(sl, 2.5, 5.05, 3.0, 5.05, color=MUTE, w=1.4, arrow=True)
marks = ['✓', '✗', '✗', '✓', '✗', '✗', '✗', '✗']
for k, mk in enumerate(marks):
    yy = 3.55 + k * 0.42
    chip(sl, 3.1, yy, 2.0, 0.32, fill=WHITE, line=LINE, radius=0.25)
    tb(sl, 5.2, yy - 0.05, 0.5, 0.4, [[R(mk, c=(TEAL if mk == '✓' else RED), s=12.5, b=True)]])
tb(sl, 3.1, 7.05, 2.6, 0.4, [[R('n = 8 rollouts', c=MUTE, s=12.5)]])
seg(sl, 5.8, 5.05, 6.3, 5.05, color=MUTE, w=1.4, arrow=True)
chip(sl, 6.4, 4.5, 1.7, 1.1, runs_=[M('P', c=INK, s=17), M('i', c=INK, s=12, sub=True), R(' = 0.25', c=INK, s=15, b=True)],
     fill=PANEL, line=LINE, radius=0.15)
seg(sl, 8.2, 5.05, 8.7, 5.05, color=MUTE, w=1.4, arrow=True)
panel(sl, 8.8, 4.3, 2.3, 1.5, fill=WHITE, line=SLATE, lw=1.6, shape=MSO_SHAPE.DIAMOND)
tb(sl, 8.8, 4.72, 2.3, 0.7, [[M('P', c=INK, s=15), M('i', c=INK, s=11, sub=True), M(' ≤ γ ?', c=INK, s=15)]],
   align=PP_ALIGN.CENTER)
# two branches
seg(sl, 11.2, 5.05, 12.3, 3.7, color=ORANGE, w=2.0, arrow=True)
tb(sl, 11.05, 3.55, 1.5, 0.4, [[R('yes · τ★ 존재', c=ORANGE, s=12.5, b=True)]])
panel(sl, 12.5, 2.95, 6.4, 1.55, fill=ORANGE_BG, line=ORANGE, lw=1.5)
tb(sl, 12.9, 3.12, 5.7, 1.3,
   [[R('SFT branch — ', b=True, c=ORANGE, s=17), R('verified trajectory ', c=SLATE, s=17), M('τ', c=SLATE, s=17), M('★', c=SLATE, s=12, sup=True)],
    [R('prompt당 1행 · 실패 영역에 supervised 신호', c=SLATE, s=14)]])
seg(sl, 11.2, 5.05, 12.3, 6.5, color=BLUE, w=2.0, arrow=True)
tb(sl, 11.35, 6.35, 1.2, 0.4, [[R('no', c=BLUE, s=12.5, b=True)]])
panel(sl, 12.5, 5.85, 6.4, 1.55, fill=BLUE_BG, line=BLUE, lw=1.5)
tb(sl, 12.9, 6.02, 5.7, 1.3,
   [[R('RL branch — ', b=True, c=BLUE, s=17), R('rollout group 그대로', c=SLATE, s=17)],
    [R('prompt당 n행 · group-relative advantage', c=SLATE, s=14)]])
seg(sl, 15.7, 4.5, 15.7, 5.85, color=ORANGE, w=1.4, dash='dash', arrow=True)
tb(sl, 15.9, 4.95, 3.0, 0.5, [[R('τ★ 없으면 RL로 fallback', c=MUTE, s=12.5)]])
band(sl, 0.9, 8.5, 18.2, 1.35,
     [[R('결정의 단위 = 신호의 단위.  ', b=True, c=TEAL, s=18),
       R('supervised 신호(τ★)도 group advantage도 prompt에 묶여 있다 — routing이 같은 grain일 때 두 branch가 정합적으로 대응한다.',
         c=SLATE, s=17)]])

# ================= S10 · 하나의 interface =================
sl = add_slide('방법', '두 branch는 하나의 learner interface로 내려온다')
panel(sl, 0.9, 2.1, 8.85, 5.4)
tb(sl, 1.35, 2.4, 8.0, 0.5, [[R('verified trajectory는 문자열이 아니라 상호작용 기록', b=True, c=INK, s=18)]])
trace = [('instruction', PANEL, SLATE), ('obs', GEN, SLATE), ('reason·action', ORANGE, WHITE),
         ('obs', GEN, SLATE), ('reason·action', ORANGE, WHITE), ('종료', ORANGE, WHITE)]
xx = 1.35
for t, f, tc in trace:
    ww = 0.55 + 0.068 * len(t)
    chip(sl, xx, 3.5, ww, 0.62, t, fill=f, tcolor=tc, size=11.5, line=LINE if f in (PANEL, GEN) else None, radius=0.25)
    xx += ww + 0.12
tb(sl, 1.35, 4.45, 8.0, 0.4,
   [[R('학습 신호는 assistant token에만 — ', c=SLATE, s=15), R('response mask', c=ORANGE, s=15, b=True),
     R('로 학습 구간을 표시', c=SLATE, s=15)]])
tb(sl, 1.35, 5.45, 8.0, 1.6,
   [[R('observation은 다음 action을 조건화하는 context일 뿐, 학습 대상이 아니다.', c=SLATE, s=15)],
    [R('materialize되면 RL row와 같은 형식 — context + response + response mask.', c=SLATE, s=15)]])
panel(sl, 10.25, 2.1, 8.85, 5.4)
tb(sl, 10.7, 2.4, 8.0, 0.5, [[R('row 비대칭 — 같은 prompt가 RL은 n행, SFT는 1행', b=True, c=INK, s=18)]])
for k in range(4):
    chip(sl, 10.7, 3.4 + k * 0.5, 2.6, 0.36, fill=BLUE, radius=0.22)
tb(sl, 10.7, 5.5, 2.8, 0.4, [[R('RL — n행', c=BLUE, s=13.5, b=True)]])
chip(sl, 14.2, 3.4, 2.6, 0.36, fill=ORANGE, radius=0.22)
tb(sl, 14.2, 3.9, 3.2, 0.4, [[R('SFT — 1행 (τ★)', c=ORANGE, s=13.5, b=True)]])
tb(sl, 10.7, 6.0, 8.1, 1.3,
   [[R('두 층의 identity로 해소 — ', c=SLATE, s=15, b=True),
     R('출처를 가리키는 prompt identity, 계산 단위를 가리키는 advantage-group identity.', c=SLATE, s=15)]])
band(sl, 0.9, 8.15, 18.2, 1.25,
     [[R('branch를 가르는 것은 별도의 학습 경로가 아니라, batch가 실어 나르는 정보다.', b=True, c=TEAL, s=19)]])

# ================= S11 · Q2a advantage 통일 =================
sl = add_slide('방법 · Q2a', 'Advantage의 통일 — singleton group과 β pseudo-reward')
panel(sl, 0.9, 2.1, 8.85, 4.7, fill=BLUE_BG, line=BLUE, lw=1.5)
tb(sl, 1.35, 2.4, 8.0, 0.5, [[R('RL — rollout group이 advantage의 단위', b=True, c=BLUE, s=18)]])
rs = ['r = 1', 'r = 0', 'r = 1', 'r = 0']
for k, rr in enumerate(rs):
    chip(sl, 1.35 + k * 1.85, 3.35, 1.6, 0.6, rr, fill=WHITE, tcolor=SLATE, size=14, line=LINE, radius=0.2)
tb(sl, 1.35, 4.25, 7.9, 0.5, [[R('… n개 rollout이 한 group', c=MUTE, s=13.5)]])
tb(sl, 1.35, 4.95, 7.9, 1.4,
   [[M('A', c=INK, s=18), M('i,j', c=INK, s=12, sub=True), M(' = r', c=INK, s=18), M('i,j', c=INK, s=12, sub=True),
     M(' − mean(group)', c=INK, s=18)],
    [R('group 안의 상대 성과가 신호가 된다', c=SLATE, s=14.5)]])
panel(sl, 10.25, 2.1, 8.85, 4.7, fill=ORANGE_BG, line=ORANGE, lw=1.5)
tb(sl, 10.7, 2.4, 8.0, 0.5, [[R('SFT — verified trajectory가 혼자 singleton group', b=True, c=ORANGE, s=18)]])
chip(sl, 10.7, 3.35, 1.6, 0.6, runs_=[M('τ', c=SLATE, s=15), M('★', c=SLATE, s=11, sup=True)],
     fill=WHITE, line=LINE, radius=0.2)
tb(sl, 12.5, 3.42, 6.3, 0.5, [[R('+ pseudo reward ', c=SLATE, s=15), M('β', c=ORANGE, s=17, b=True)]])
tb(sl, 10.7, 4.25, 8.0, 0.5, [[R('비교할 다른 rollout이 없으므로, β가 그대로 advantage가 된다  (std 정규화 없는 group 경로)', c=MUTE, s=13)]])
tb(sl, 10.7, 4.95, 8.0, 1.4,
   [[M('A', c=INK, s=18), M('r,t', c=INK, s=12, sub=True), M(' = β · m', c=INK, s=18), M('r,t', c=INK, s=12, sub=True),
     R('   (supervised token mask 위에 균일)', c=SLATE, s=14.5)],
    [R('positive supervised 신호가 group 경로에서 그대로 나온다', c=SLATE, s=14.5)]])
band(sl, 0.9, 7.35, 18.2, 1.6,
     [[R('새 estimator를 추가하지 않는다.  ', b=True, c=TEAL, s=20),
       R('singleton group과 β라는 “입력”만으로, RL이 쓰는 group-relative advantage 경로가 supervised 신호를 생성한다.',
         c=SLATE, s=18)]])

# ================= S12 · Q2b reference 통일 =================
sl = add_slide('방법 · Q2b', 'Reference의 통일 — self-detach가 SFT를 복원한다')
panel(sl, 0.9, 2.05, 11.2, 6.1)
tb(sl, 1.35, 2.35, 10.4, 0.45, [[R('effective old log-probability — branch indicator ', c=SLATE, s=15),
                                 M('z', c=SLATE, s=15), M('r', c=SLATE, s=11, sub=True),
                                 R(' (SFT = 1)', c=MUTE, s=13.5)]])
tb(sl, 1.35, 2.95, 10.4, 0.7,
   [[M('ℓ̃', s=22), M('old', s=14, sup=True), M('  =  (1 − z', s=22), M('r', s=14, sub=True),
     M(') · ℓ', s=22), M('rollout', s=14, sup=True), M('  +  z', s=22), M('r', s=14, sub=True),
     M(' · stopgrad( ℓ', s=22), M('θ', s=14, sub=True), M(' )', s=22)]], color=INK)
tb(sl, 1.35, 3.95, 10.4, 0.7,
   [[M('ρ  =  exp( ℓ', s=22), M('θ', s=14, sub=True), M(' − ℓ̃', s=22), M('old', s=14, sup=True), M(' )', s=22),
     R('     → 같은 clipped PPO objective 통과', c=SLATE, s=15)]], color=INK)
seg(sl, 1.35, 4.95, 11.6, 4.95, color=LINE, w=1.0)
tb(sl, 1.35, 5.15, 10.4, 0.55,
   [[R('SFT 행: forward에서 ', c=SLATE, s=16.5), M('ρ ≡ 1', c=INK, s=18, b=True),
     R(' — clip은 항등으로 비활성.  그러나 gradient는 살아 있다:', c=SLATE, s=16.5)]])
tb(sl, 1.35, 5.9, 10.4, 0.75,
   [[M('∇', s=21), M('θ', s=13, sub=True), M(' ρ  =  ∇', s=21), M('θ', s=13, sub=True),
     M(' log π', s=21), M('θ', s=13, sub=True),
     M('      ⇒      update  ∝  A · (1/π', s=21), M('θ', s=13, sub=True),
     M(') ∇', s=21), M('θ', s=13, sub=True), M(' π', s=21), M('θ', s=13, sub=True)]], color=INK)
tb(sl, 1.35, 6.85, 10.4, 1.2,
   [[R('stopgrad는 forward 값만 고정하고 미분에는 참여하지 않는다 — 분모가 ', c=SLATE, s=15),
     M('π', c=INK, s=16), M('θ', c=INK, s=11, sub=True),
     R('인 SFT의 gradient가 정확히 남는다.', c=SLATE, s=15)],
    [R('직관 — ', c=TEAL, s=14, b=True),
     R('stopgrad로 감싼 값은 “읽을 때는 상수, 미분할 때는 없는 것”처럼 행동한다.', c=MUTE, s=14)]])
# right recap
panel(sl, 12.55, 2.05, 6.55, 6.1, fill=PANEL, line=LINE)
tb(sl, 12.95, 2.35, 5.8, 0.5, [[R('사전지식 ②의 표, 회수', b=True, c=TEAL, s=17)]])
recap = [('data source', 'rollout ↔ τ★', '라우팅 (Q1)'),
         ('advantage', 'group ↔ β singleton', 'Q2a'),
         ('reference', 'rollout logπ ↔ self-detach', '이 슬라이드'),
         ('stabilization', 'clip 활성 ↔ 항등 비활성', '이 슬라이드')]
yy = 3.05
for a, b_, c_ in recap:
    tb(sl, 12.95, yy, 5.9, 0.75,
       [[R('✓ ', c=TEAL, s=15, b=True), R(a, c=INK, s=15, b=True), R('  ' + b_, c=SLATE, s=14)],
        [R('     ' + c_, c=MUTE, s=12)]], space_after=2)
    yy += 1.02
tb(sl, 12.95, 7.3, 5.9, 0.7, [[R('네 슬롯 전부가 하나의 ratio 경로 위에서 채워졌다.', c=SLATE, s=14.5, b=True)]])
band(sl, 0.9, 8.45, 18.2, 1.7,
     [[R('별도의 cross-entropy loss 없이 — ', b=True, c=TEAL, s=18.5),
       R('self-detach 한 번으로 SFT가 RL과 같은 PPO ratio 경로 위에서 복원된다.', c=SLATE, s=18.5)],
      [R('말로 하면: 안전장치(clip)에는 “변화 없음(ρ=1)”으로 보이게 하고, 학습 신호로는 “이 정답 token의 확률을 올려라”만 남긴다.',
         c=MUTE, s=15)]])

# ================= S13 · Q3 provenance =================
sl = add_slide('방법 · Q3', 'Provenance — validity 기제는 RL 행에만 의미를 갖는다')
panel(sl, 0.9, 2.1, 18.2, 2.15, fill=BLUE_BG, line=BLUE, lw=1.5)
tb(sl, 1.35, 2.35, 3.0, 0.5, [[R('RL 행', b=True, c=BLUE, s=19)]])
tb(sl, 1.35, 2.95, 8.6, 1.1,
   [[R('rollout logπ ✓ · 생성 version g ✓', c=INK, s=16, b=True)],
    [R('비동기 rollout data — provenance가 있다', c=SLATE, s=14.5)]])
chip(sl, 10.6, 2.7, 3.9, 0.95, 'correction 적용', fill=BLUE, size=15.5, radius=0.16)
chip(sl, 14.7, 2.7, 3.9, 0.95, 'staleness 필터 대상', fill=BLUE, size=15.5, radius=0.16)
panel(sl, 0.9, 4.5, 18.2, 2.15, fill=ORANGE_BG, line=ORANGE, lw=1.5)
tb(sl, 1.35, 4.75, 3.2, 0.5, [[R('SFT 행', b=True, c=ORANGE, s=19)]])
tb(sl, 1.35, 5.35, 8.6, 1.1,
   [[R('rollout logπ = placeholder · 생성 version 없음', c=INK, s=16, b=True)],
    [R('고정된 verified target — provenance가 없다', c=SLATE, s=14.5)]])
chip(sl, 10.6, 5.1, 8.0, 0.95, 'identity 통과 — 보정도 필터도 적용되지 않는다', fill=ORANGE, size=15.5, radius=0.16)
tb(sl, 0.9, 7.0, 18.2, 0.9,
   [[R('같은 모양의 tensor, 다른 의미 — ', b=True, c=INK, s=16.5),
     R('rollout logπ는 RL 행에서는 reference이자 mismatch의 증거, SFT 행에서는 batch schema를 맞추는 값일 뿐이다. '
       'validity 기제는 placeholder를 증거로 해석하지 않는다.', c=SLATE, s=16.5)]])
band(sl, 0.9, 8.3, 18.2, 1.45,
     [[R('사전지식 ①의 회수 — ', b=True, c=TEAL, s=18.5),
       R('correction과 staleness는 provenance에 의존하는 기제였다. provenance가 있는 행에만 적용하면, 도입의 충돌은 정확히 이 지점에서 해소된다.',
         c=SLATE, s=17.5)]])

# ================= S14 · Q4 실행 단위 =================
sl = add_slide('방법 · Q4', '실행의 단위와 의미의 단위')
# semantic lane
tb(sl, 0.9, 2.25, 3.4, 0.5, [[R('semantic unit', b=True, c=TEAL, s=17)], [R('prompt group', c=MUTE, s=13)]])
chip(sl, 4.5, 2.25, 2.5, 0.9, 'prompt group', fill=TEAL, size=14.5, radius=0.16)
seg(sl, 7.1, 2.7, 12.4, 2.7, color=MUTE, w=1.4, dash='dash')
tb(sl, 7.3, 2.25, 5.0, 0.4, [[R('실행 동안 의미론에는 보이지 않음', c=MUTE, s=12.5)]], align=PP_ALIGN.CENTER)
chip(sl, 12.5, 2.25, 2.9, 0.9, 'group 복원', fill=TEAL, size=14.5, radius=0.16)
seg(sl, 15.5, 2.7, 16.1, 2.7, color=MUTE, w=1.4, arrow=True)
chip(sl, 16.2, 2.25, 2.2, 0.9, 'Routing Gate', fill=WHITE, line=SLATE, tcolor=INK, size=14.5, radius=0.16)
# execution lane
tb(sl, 0.9, 5.25, 3.4, 0.7, [[R('execution unit', b=True, c=INK, s=17)], [R('trajectory attempt', c=MUTE, s=13)]])
alens = [2.0, 2.8, 4.6, 2.3]
for k, alw in enumerate(alens):
    yy = 4.55 + k * 0.72
    f = RED_BG if k == 2 else GEN
    bd = RED if k == 2 else None
    chip(sl, 5.4, yy, alw, 0.5, f'attempt {k+1}', fill=f, tcolor=(RED if k == 2 else SLATE), size=12.5,
         line=bd, radius=0.2)
    seg(sl, 5.4 + alw + 0.1, yy + 0.25, 12.6, yy + 0.25, color=LINE, w=1.0, dash='sysDot')
tb(sl, 10.3, 6.05, 2.2, 0.4, [[R('long-tail', c=RED, s=12.5, b=True)]])
seg(sl, 5.0, 3.15, 5.0, 4.6, color=MUTE, w=1.4, arrow=True)
tb(sl, 5.25, 3.65, 3.4, 0.4, [[R('n개로 분해, 독립 실행', c=MUTE, s=12.5)]])
seg(sl, 12.8, 6.3, 13.6, 3.15, color=MUTE, w=1.4, arrow=True)
tb(sl, 13.95, 4.85, 4.6, 0.4, [[R('완료된 attempt를 accumulator가 다시 모음', c=MUTE, s=12.5)]])
tb(sl, 0.9, 7.5, 18.2, 0.5,
   [[R('분할은 실행에만 존재한다 — gate가 success profile을 보는 시점에는 이미 사라져 있다.', b=True, c=INK, s=17)]])
notes = [
    ('long-tail 해소', '긴 attempt가 group 전체의 완료를 붙잡지 않는다 — 실행 효율은 attempt 단위에서 나온다.'),
    ('backpressure도 attempt 단위', '완료 group + 완성 대기 attempt를 합산해 계량 — 부분 완료 backlog가 보이지 않는 상태로 남지 않는다.'),
    ('원리(A.1)와 같은 사상', '의미론이 보는 세계를 실행의 사정으로부터 격리한다 — routing·contract·objective는 이 절의 존재를 몰라도 성립한다.'),
]
for i, (t, d) in enumerate(notes):
    x = 0.9 + i * 6.2
    panel(sl, x, 8.25, 5.8, 1.85)
    tb(sl, x + 0.35, 8.45, 5.15, 0.45, [[R(t, b=True, c=TEAL, s=15)]])
    tb(sl, x + 0.35, 8.95, 5.15, 1.05, [[R(d, c=SLATE, s=13.5)]])

# ================= S15 · 실험 셋업 =================
sl = add_slide('실험', '셋업 — 통제된 math RLVR testbed')
band(sl, 0.9, 2.0, 18.2, 1.0,
     [[R('본 도메인은 computer-use post-training — ', b=True, c=SLATE, s=17),
       R('설계 검증은 신호가 통제된 math RLVR testbed에서 수행한다.', c=SLATE, s=17)]],
     fill=PANEL, border=LINE)
cells = [
    ('Task · Data', 'OpenR1 math prompts + verified trajectory τ★ (optional coverage — 없는 prompt는 RL fallback)'),
    ('Model', 'Qwen2.5-Math-1.5B (주 실험) · 7B 확인'),
    ('Rollout', '비동기 stream · prompt당 group n = 8 · rule-based verifier'),
    ('평가', 'held-out math 벤치마크 · mean@8'),
    ('비교 anchor', 'D0 — 분리 없는 coupled 구성 (all-off 코너)'),
    ('구현', 'verl 기반 비동기 스택 위에 구현 — 별도 학습 루프 없이 단일 learner'),
]
for i, (t, d) in enumerate(cells):
    x = 0.9 + (i % 3) * 6.2
    y = 3.35 + (i // 3) * 2.55
    panel(sl, x, y, 5.8, 2.3)
    tb(sl, x + 0.35, y + 0.25, 5.1, 0.45, [[R(t, b=True, c=INK, s=17)]])
    tb(sl, x + 0.35, y + 0.85, 5.1, 1.35, [[R(d, c=SLATE, s=14.5)]])
tb(sl, 0.9, 8.85, 18.2, 0.6,
   [[R('두 branch의 기여는 row가 아니라 prompt group 단위로 계측한다 — RL은 한 prompt가 n행으로 펼쳐지므로.',
       c=MUTE, s=14.5)]])

# ================= S16 · 메인 결과 (placeholder) =================
sl = add_slide('실험', '메인 결과 — coupled anchor 대비')
for i, (t, d) in enumerate([
    ('학습 곡선 — reward · 안정성', 'D0 (coupled) vs Async-HPT 구성 — 학습 진행에 따른 보상과 안정성 지표'),
    ('벤치마크 — mean@8', 'held-out math 벤치마크 최종 성능 표'),
]):
    x = 0.9 + i * 9.35
    sp = panel(sl, x, 2.1, 8.85, 6.2, fill=WHITE, line=MUTE, lw=1.4)
    dash_border(sp)
    tb(sl, x + 0.45, 2.4, 8.0, 0.5, [[R(t, b=True, c=INK, s=18)]])
    tb(sl, x + 0.45, 4.7, 7.95, 0.8, [[R('〔 런 결과로 채움 〕', c=MUTE, s=20, b=True)]], align=PP_ALIGN.CENTER)
    tb(sl, x + 0.45, 7.35, 7.95, 0.7, [[R(d, c=MUTE, s=13.5)]])
band(sl, 0.9, 8.75, 18.2, 1.15,
     [[R('읽는 법 — ', b=True, c=TEAL, s=17),
       R('같은 데이터·같은 라우팅에서 비동기 처리와 objective 구성만 바꾼 비교. 의미 보존 설계의 효과를 직접 읽는다.',
         c=SLATE, s=17)]])

# ================= S17 · Ablation 2×2 =================
sl = add_slide('실험', 'Ablation — 2×2: 비동기 보정의 분리 × objective 안정화')
tb(sl, 4.9, 2.1, 12.0, 0.5,
   [[R('C2 — objective 안정화', b=True, c=INK, s=16),
     R('   “안정화 방식은 독립적으로 기여하는가?”', c=MUTE, s=14)]])
chip(sl, 4.9, 2.7, 5.9, 0.7, runs_=[R('vanilla clip', c=INK, s=14.5, b=True),
                                    R(' — 표준 PPO clip', c=MUTE, s=12.5)],
     fill=PANEL, line=LINE, radius=0.15)
chip(sl, 11.0, 2.7, 5.9, 0.7, runs_=[R('CISPO', c=INK, s=14.5, b=True),
                                     R(' — clip 대신 IS weight 절단으로 안정화', c=MUTE, s=12.5)],
     fill=PANEL, line=LINE, radius=0.15)
tb(sl, 0.9, 2.1, 3.7, 1.4,
   [[R('C1 — 비동기 보정의 분리', b=True, c=INK, s=16)],
    [R('“staleness 보정을 objective의 기준점(anchor logπ)에서 분리해야 하는가?”', c=MUTE, s=13)]])
chip(sl, 0.9, 4.25, 3.7, 0.7, 'coupled', fill=PANEL, tcolor=INK, size=15, line=LINE, radius=0.15)
chip(sl, 0.9, 6.4, 3.7, 0.7, 'decoupled', fill=PANEL, tcolor=INK, size=15, line=LINE, radius=0.15)
for x, y, t, d, f, bd, tc in [
    (4.9, 3.6, 'D0 계열', 'all-off 앵커 — 분리도 안정화도 없음', PANEL, LINE, INK),
    (11.0, 3.6, 'M − dec', 'coupled + CISPO — 분리 없이 안정화만', BLUE_BG, BLUE, BLUE),
    (4.9, 5.75, 'M − cispo', 'decoupled + vanilla — 분리만', BLUE_BG, BLUE, BLUE),
    (11.0, 5.75, 'M (풀스택)', 'decoupled + CISPO — 제안 구성', TEAL_BG, TEAL, TEAL),
]:
    panel(sl, x, y, 5.9, 1.95, fill=f, line=bd, lw=1.5)
    tb(sl, x + 0.4, y + 0.25, 5.1, 0.5, [[R(t, b=True, c=tc, s=19)]])
    tb(sl, x + 0.4, y + 0.9, 5.1, 0.9, [[R(d, c=SLATE, s=14.5)]])
tb(sl, 17.2, 3.6, 2.0, 4.1,
   [[R('leave-one-out', c=MUTE, s=13, b=True)],
    [R('각 축을 하나씩 빼서 기여를 분리', c=MUTE, s=13)]])
band(sl, 0.9, 8.9, 18.2, 1.0,
     [[R('집계층 ablation(A1)은 별도 축으로 통제 — 전 격자점에서 라우팅·스케일 인자는 고정.', c=SLATE, s=15.5)]],
     fill=PANEL, border=LINE)

# ================= S18 · 분석 =================
sl = add_slide('실험', '분석 — 학습 동역학에서 무엇을 볼 것인가')
qs = [
    ('신호 배분', 'SFT로 route되는 비율이 학습이 진행되며 줄어드는가 — 실패 영역이 축소된다는 직접 증거. prompt group 단위로 계측.'),
    ('안정성', '두 branch가 섞인 batch에서 entropy와 clip 통계가 안정한가 — supervised 신호가 RL 신호를 교란하지 않는가.'),
    ('길이 동역학', 'response length가 보상 개선 없이 늘어나지 않는가 — 비종료·장황화는 실패로 계측되어야 한다.'),
]
for i, (t, d) in enumerate(qs):
    x = 0.9 + i * 6.2
    panel(sl, x, 2.1, 5.8, 3.1)
    tb(sl, x + 0.35, 2.4, 5.1, 0.5, [[R('%d) %s' % (i + 1, t), b=True, c=INK, s=18)]])
    tb(sl, x + 0.35, 3.05, 5.1, 2.0, [[R(d, c=SLATE, s=14.5)]])
    sp = panel(sl, x, 5.5, 5.8, 3.4, fill=WHITE, line=MUTE, lw=1.2)
    dash_border(sp)
    tb(sl, x, 6.9, 5.8, 0.6, [[R('〔 런 결과로 채움 〕', c=MUTE, s=16, b=True)]], align=PP_ALIGN.CENTER)
tb(sl, 0.9, 9.15, 18.2, 0.6,
   [[R('세 지표 모두 branch별로 분리해 읽는다 — mixed batch의 평균은 두 branch의 병리를 서로 가릴 수 있다.',
       c=MUTE, s=14.5)]])

# ================= S19 · 결론 =================
sl = add_slide('결론', '의미를 보존하는 결합')
concl = [
    ('문제', '효율을 위한 asynchrony와 신호를 위한 hybrid estimator는 learner의 경계에서 충돌한다 — 비동기 기제의 전제(모든 sample = rollout)를 verified trajectory가 깨기 때문에.'),
    ('원리', 'transport와 semantics를 분리한다 — 경험은 하나의 stream으로 흐르되, branch의 의미는 learner boundary에서만 실행된다.'),
    ('결과', '라우팅·advantage·reference·validity·실행의 다섯 지점에서 계약을 세우면, asynchrony는 효율을, hybrid estimator는 신호를 서로의 전제를 훼손하지 않고 제공한다.'),
]
yy = 2.15
for t, d in concl:
    chip(sl, 0.9, yy + 0.08, 1.5, 0.6, t, fill=TEAL, size=15, radius=0.25)
    tb(sl, 2.75, yy, 16.3, 1.3, [[R(d, c=SLATE, s=17)]])
    yy += 1.55
panel(sl, 0.9, 7.0, 8.85, 2.6)
tb(sl, 1.35, 7.25, 8.0, 0.5, [[R('한계', b=True, c=INK, s=17)]])
tb(sl, 1.35, 7.85, 7.95, 1.6,
   [[R('· 검증은 math RLVR testbed — 본 도메인(computer-use)의 규모는 남은 일', c=SLATE, s=14.5)],
    [R('· version 기반 staleness 필터는 설계로 지원하되 본 실험에서는 비활성', c=SLATE, s=14.5)]])
panel(sl, 10.25, 7.0, 8.85, 2.6)
tb(sl, 10.7, 7.25, 8.0, 0.5, [[R('향후', b=True, c=INK, s=17)]])
tb(sl, 10.7, 7.85, 7.95, 1.6,
   [[R('· computer-use agent post-training으로의 본 적용 — 긴 상호작용 trace, 희박한 검증 보상', c=SLATE, s=14.5)],
    [R('· τ★ coverage·라우팅 임계의 적응적 운용', c=SLATE, s=14.5)]])
tb(sl, 0.9, 9.9, 18.2, 0.6, [[R('감사합니다 — Questions?', b=True, c=INK, s=20)]], align=PP_ALIGN.CENTER)

# ================= S20 · Backup divider =================
sl = add_slide(footer=False)
tb(sl, 0.9, 4.9, 18.2, 1.0, [[R('Backup', b=True, c=INK, s=44)]], align=PP_ALIGN.CENTER)
tb(sl, 0.9, 6.05, 18.2, 0.6,
   [[R('reduction과 계수 · self-detach 유도 · 라우팅 회귀 사례', c=MUTE, s=17)]], align=PP_ALIGN.CENTER)

# ================= S21 · Backup: reduction과 계수 =================
sl = add_slide('Backup', 'Reduction과 계수 — α와 β는 다른 자리에서 작동한다')
tb(sl, 0.9, 2.05, 18.2, 0.9,
   [[R('원 HPT는 ', c=SLATE, s=17), M('α L', c=INK, s=18), M('RL', c=INK, s=12, sub=True),
     M(' + β L', c=INK, s=18), M('SFT', c=INK, s=12, sub=True),
     R('의 loss mixing — Async-HPT에서는 이 선택이 route 단위 branch 선택으로 내려온다.', c=SLATE, s=17)]])
rows21 = [
    ('sequence weight', [M('w', s=17), M('r', s=11, sub=True), M(' = α/n', s=17)],
     [M('w', s=17), M('r', s=11, sub=True), M(' = 1', s=17)],
     '한 prompt의 총 기여를 row 수와 무관하게 한 단위로'),
    ('length divisor', [M('d', s=17), M('r', s=11, sub=True), M(' = T', s=17), M('rollout', s=11, sup=True)],
     [M('d', s=17), M('r', s=11, sub=True), M(' = Σ m', s=17), M('r,t', s=11, sub=True)],
     'RL은 materialize 시점 response width, SFT는 supervised token 수'),
    ('scale의 자리', [R('α — prompt-level RL 기여')], [R('β — pseudo reward (advantage 크기)')],
     '같은 계수가 branch에 따라 다른 위치에서 작동'),
]
ty = 3.3
chip(sl, 0.9, ty, 6.7, 0.7, '항목', fill=INK, size=15, radius=0.12)
chip(sl, 7.7, ty, 5.2, 0.7, 'RL branch', fill=BLUE, size=15, radius=0.12)
chip(sl, 13.0, ty, 5.2, 0.7, 'SFT branch', fill=ORANGE, size=15, radius=0.12)
for i, (name, rl, sft, note) in enumerate(rows21):
    yy = ty + 0.8 + i * 1.45
    panel(sl, 0.9, yy, 6.7, 1.3, fill=PANEL, line=LINE, radius=0.1)
    tb(sl, 1.25, yy + 0.15, 6.1, 1.05,
       [[R(name, c=INK, s=15.5, b=True)], [R(note, c=MUTE, s=12)]], space_after=2)
    panel(sl, 7.7, yy, 5.2, 1.3, fill=BLUE_BG, line=None, radius=0.1)
    tb(sl, 7.7, yy, 5.2, 1.3, [rl], size=15.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    panel(sl, 13.0, yy, 5.2, 1.3, fill=ORANGE_BG, line=None, radius=0.1)
    tb(sl, 13.0, yy, 5.2, 1.3, [sft], size=15.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
band(sl, 0.9, 8.6, 18.2, 1.1,
     [[R('구현 확정 — ', b=True, c=TEAL, s=16),
       R('reduction은 branch-blind(token-count 기반)로 수렴, α = 1 고정. β만 pseudo reward scale로 남는다 (DR-001).',
         c=SLATE, s=16)]])

# ================= S22 · Backup: self-detach 유도 =================
sl = add_slide('Backup', 'Self-detach gradient 유도')
lines22 = [
    [R('SFT 행에서 ', c=SLATE, s=17), M('ℓ̃', c=INK, s=19), M('old', c=INK, s=12, sup=True),
     M(' = stopgrad(ℓ', c=INK, s=19), M('θ', c=INK, s=12, sub=True), M(')', c=INK, s=19),
     R('이므로 forward 값은 ', c=SLATE, s=17), M('ρ = 1', c=INK, s=19, b=True)],
    [M('∇', s=20), M('θ', s=13, sub=True), M(' ρ  =  ρ · ∇', s=20), M('θ', s=13, sub=True),
     M(' ℓ', s=20), M('θ', s=13, sub=True), M('  =  ∇', s=20), M('θ', s=13, sub=True),
     M(' log π', s=20), M('θ', s=13, sub=True), M('(a | h)  =  (1/π', s=20), M('θ', s=13, sub=True),
     M(') ∇', s=20), M('θ', s=13, sub=True), M(' π', s=20), M('θ', s=13, sub=True)],
    [R('token loss가 ', c=SLATE, s=17), M('−A · ρ', c=INK, s=19),
     R(' 형태일 때 update direction은', c=SLATE, s=17)],
    [M('A · (1/π', s=20), M('θ', s=13, sub=True), M(') ∇', s=20), M('θ', s=13, sub=True),
     M(' π', s=20), M('θ', s=13, sub=True),
     R('   —  reference denominator가 현재 policy이고 advantage가 상수인, unified estimator의 SFT 형태와 일치',
       c=SLATE, s=15.5)],
]
yy = 2.4
for ln in lines22:
    tb(sl, 1.3, yy, 17.4, 0.85, [ln], color=INK)
    yy += 1.25
band(sl, 0.9, 7.6, 18.2, 1.5,
     [[R('핵심 — ', b=True, c=TEAL, s=17),
       R('stopgrad는 forward 값만 고정한다. clip·IS·advantage 경로를 전부 공유하면서도, gradient 층위에서는 정확히 SFT가 남는다.',
         c=SLATE, s=17)]])

# ================= S23 · Backup: 라우팅 회귀 사례 =================
sl = add_slide('Backup', '라우팅은 단일 실패 지점이다 — 한 회귀 사례')
tb(sl, 0.9, 2.1, 18.2, 0.9,
   [[R('gate는 성공 신호 하나로 전체 데이터 흐름을 가른다. 성공 신호를 잘못 읽으면 — 예컨대 보상 tensor의 다른 위치를 참조하면 — ',
       c=SLATE, s=17),
     R('전 prompt가 실패로 읽혀 학습이 사실상 pure-SFT로 붕괴한다.', c=RED, s=17, b=True)]])
obs = [
    ('증상', 'RL branch 비율 0 근접 · entropy와 response length의 동반 이상 — 원인은 objective가 아니라 라우팅 입력.'),
    ('교훈 ①', '라우팅 입력(성공률 통계)과 branch 비율은 학습 지표와 같은 등급의 상시 모니터링 대상이다.'),
    ('교훈 ②', '비종료·잘림 응답은 실패로 계측되어야 한다 — 관대한 채점은 라우팅과 보상 양쪽을 함께 왜곡한다.'),
]
yy = 3.55
for t, d in obs:
    chip(sl, 0.9, yy + 0.05, 1.9, 0.6, t, fill=INK, size=14, radius=0.22)
    tb(sl, 3.1, yy, 16.0, 1.0, [[R(d, c=SLATE, s=16)]])
    yy += 1.3
band(sl, 0.9, 7.8, 18.2, 1.3,
     [[R('설계 함의 — ', b=True, c=TEAL, s=17),
       R('transport/semantics 분리는 이런 실패를 국소화한다: 라우팅 결함은 gate 한 곳의 문제로 남고, objective·validity 계약은 그대로 성립한다.',
         c=SLATE, s=17)]])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'AsyncHPT_Slides.pptx')
prs.save(out)
print('saved:', out, '| slides:', len(prs.slides.__iter__.__self__._sldIdLst))
